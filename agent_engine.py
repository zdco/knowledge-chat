"""全能 Agent 引擎：知识域动态加载 + 工具定义 + 执行 + 多 SDK 流式调用"""
import os
import re
import json
import subprocess
import glob as glob_mod
import time
import contextvars
import urllib.request
import urllib.error
import threading
import logging

import yaml
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# Office 文件扩展名集合
_OFFICE_EXTS = {'.xlsx', '.xls', '.docx', '.pptx', '.pdf'}

# ── 加载配置 ──────────────────────────────────────────────

_DIR = os.path.dirname(__file__)
_CONFIG_PATH = os.path.join(_DIR, "config.yaml")


def _load_config() -> dict:
    with open(_CONFIG_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


CONFIG = _load_config()

BASE_URL = os.environ.get("AI_CHAT_BASE_URL", CONFIG["api"]["base_url"])
API_KEY = os.environ.get("AI_CHAT_API_KEY", CONFIG["api"]["api_key"])
MODEL = os.environ.get("AI_CHAT_MODEL", CONFIG["api"]["model"])
MAX_TOKENS = CONFIG["api"]["max_tokens"]
MAX_ITERATIONS = CONFIG["api"]["max_iterations"]
MAX_OUTPUT_LEN = CONFIG["tools"]["max_output_length"]
MAX_DISPLAY_LEN = CONFIG["tools"]["max_display_length"]
_COMPACT_KEEP_RECENT = CONFIG["tools"].get("compact_keep_recent", 3)
_COMPACT_MAX_LEN = CONFIG["tools"].get("compact_max_length", 500)
_COMPACT_SKIP_TOOLS = {"switch_service", "list_services", "scan_service"}

# API 格式：显式配置 > 根据实际生效的 BASE_URL 自动判断（含 /v1 用 openai，否则 anthropic）
_api_format_cfg = CONFIG["api"].get("api_format", "").lower()
if _api_format_cfg in ("openai", "anthropic"):
    API_FORMAT = _api_format_cfg
else:
    API_FORMAT = "openai" if "/v1" in BASE_URL else "anthropic"

PROJECT_ROOT = os.path.abspath(_DIR)

# 运行模式
APP_MODE = CONFIG.get("mode", "knowledge")

# ── 日志分析模式初始化 ────────────────────────────────────
_session_manager = None
_analyzer_services = {}
_analyzer_businesses = {}

if APP_MODE == "log-analyzer":
    from log_analyzer import SessionManager, load_services_config, load_businesses_config, \
        scan_service_deps, read_log_filtered, extract_log_summary, process_upload, is_image_file

    _analyzer_cfg = CONFIG.get("analyzer", {})
    _services_config_path = os.path.join(_DIR, _analyzer_cfg.get("services_config", "services.yaml"))
    _analyzer_services = load_services_config(_services_config_path)
    _analyzer_businesses = load_businesses_config(_services_config_path)
    _session_manager = SessionManager(
        session_dir=_analyzer_cfg.get("session_dir", "/data/sessions"),
        worktree_base=_analyzer_cfg.get("worktree_base", "/data/worktrees"),
        session_ttl=_analyzer_cfg.get("session_ttl", 86400),
    )
    _session_manager.cleanup_expired()
    _session_manager.start_cleanup_timer()


def _find_service(name: str) -> tuple[str, dict] | tuple[None, None]:
    """按 ID 或别名查找服务，返回 (service_id, service_config) 或 (None, None)"""
    if not _analyzer_services:
        return None, None
    # 先精确匹配 ID
    if name in _analyzer_services:
        return name, _analyzer_services[name]
    # 再匹配别名（不区分大小写）
    name_lower = name.lower()
    for sid, svc in _analyzer_services.items():
        for alias in (svc.get("aliases") or []):
            if alias.lower() == name_lower:
                return sid, svc
    return None, None


# ── Oracle Client 自动安装 ────────────────────────────────

def _ensure_oracle_client() -> str:
    """根据配置返回 Oracle Client 路径，配置为 auto 时自动下载安装。返回空字符串表示不启用。"""
    cfg = CONFIG["tools"].get("oracle_client_path", "")
    if not cfg:
        return ""
    if cfg != "auto":
        return cfg
    # auto 模式：自动下载到项目目录下
    install_dir = os.path.join(PROJECT_ROOT, "oracle_client")
    # 查找已有的 instantclient 目录
    for name in os.listdir(install_dir) if os.path.isdir(install_dir) else []:
        candidate = os.path.join(install_dir, name)
        if name.startswith("instantclient") and os.path.isdir(candidate):
            return candidate
    # 下载安装
    import platform
    import zipfile
    arch = platform.machine()
    if arch == "x86_64":
        url = "https://download.oracle.com/otn_software/linux/instantclient/1924000/instantclient-basic-linux.x64-19.24.0.0.0dbru.zip"
    elif arch == "aarch64":
        url = "https://download.oracle.com/otn_software/linux/instantclient/1924000/instantclient-basic-linux.arm64-19.24.0.0.0dbru.zip"
    else:
        logging.getLogger(__name__).warning("不支持的架构 %s，跳过 Oracle Client 自动安装", arch)
        return ""
    os.makedirs(install_dir, exist_ok=True)
    zip_path = os.path.join(install_dir, "instantclient.zip")
    log = logging.getLogger(__name__)
    try:
        if not os.path.isfile(zip_path):
            log.info("正在下载 Oracle Instant Client...")
            def _progress(block_num, block_size, total_size):
                downloaded = block_num * block_size
                if total_size > 0:
                    pct = min(100, downloaded * 100 // total_size)
                    mb_done = downloaded / 1024 / 1024
                    mb_total = total_size / 1024 / 1024
                    print(f"\r下载 Oracle Instant Client: {mb_done:.1f}/{mb_total:.1f} MB ({pct}%)", end="", flush=True)
            urllib.request.urlretrieve(url, zip_path, reporthook=_progress)
            print()  # 换行
        log.info("正在解压 Oracle Instant Client...")
        with zipfile.ZipFile(zip_path, "r") as zf:
            zf.extractall(install_dir)
        os.remove(zip_path)
    except Exception as e:
        log.error("Oracle Client 下载失败: %s", e)
        return ""
    # 查找解压后的目录
    for name in os.listdir(install_dir):
        candidate = os.path.join(install_dir, name)
        if name.startswith("instantclient") and os.path.isdir(candidate):
            logging.getLogger(__name__).info("Oracle Instant Client 已安装: %s", candidate)
            return candidate
    return ""


ORACLE_CLIENT_PATH = _ensure_oracle_client()

# 数据库密码映射，由 build_system_prompt() 填充
_DB_PASSWORDS: dict[str, str] = {}

# ── 知识域加载 ────────────────────────────────────────────

_KNOWLEDGE_DIR = os.path.join(_DIR, "knowledge")

logger = logging.getLogger(__name__)


def load_knowledge_domains() -> list[dict]:
    """扫描 knowledge/*/domain.yaml，返回知识域列表（跳过 _template）"""
    domains = []
    if not os.path.isdir(_KNOWLEDGE_DIR):
        return domains
    for dname in sorted(os.listdir(_KNOWLEDGE_DIR)):
        if dname.startswith("_") or dname.startswith("."):
            continue
        domain_dir = os.path.join(_KNOWLEDGE_DIR, dname)
        if not os.path.isdir(domain_dir):
            continue
        fpath = os.path.join(domain_dir, "domain.yaml")
        if not os.path.isfile(fpath):
            continue
        with open(fpath, "r", encoding="utf-8") as f:
            domain = yaml.safe_load(f)
        if domain and isinstance(domain, dict):
            # 解析 data_path 为绝对路径，追加到 search_paths
            data_path = domain.get("data_path")
            if data_path:
                abs_data_path = os.path.normpath(os.path.join(domain_dir, data_path))
                domain["_abs_data_path"] = abs_data_path
                # 计算相对于项目根目录的路径，追加到 search_paths
                rel_data_path = os.path.relpath(abs_data_path, PROJECT_ROOT)
                search_paths = domain.get("search_paths", [])
                if rel_data_path not in search_paths:
                    search_paths.append(rel_data_path)
                    domain["search_paths"] = search_paths

            # Confluence zip 自动转换（支持单个字符串或列表）
            confluence_zip = domain.get("confluence_zip")
            if confluence_zip:
                if isinstance(confluence_zip, str):
                    confluence_zip = [confluence_zip]
                for zip_name in confluence_zip:
                    zip_path = os.path.join(domain_dir, zip_name)
                    wiki_dir = os.path.join(
                        abs_data_path if data_path else domain_dir, "wiki"
                    )
                    if os.path.isfile(zip_path):
                        try:
                            from confluence_converter import convert_confluence_zip
                            convert_confluence_zip(zip_path, wiki_dir, domain.get("name", dname))
                        except Exception as e:
                            logging.getLogger(__name__).error(
                                "Confluence 转换失败 [%s] %s: %s", dname, zip_name, e, exc_info=True
                            )

            domains.append(domain)
    return domains


KNOWLEDGE_DOMAINS = load_knowledge_domains() if APP_MODE == "knowledge" else []

_domains_lock = threading.Lock()
request_id_var: contextvars.ContextVar[str] = contextvars.ContextVar("request_id", default="-")


def reload_domains():
    """重新加载所有知识域并重建 system prompt"""
    if APP_MODE != "knowledge":
        return
    global KNOWLEDGE_DOMAINS, SYSTEM_PROMPT
    new_domains = load_knowledge_domains()
    with _domains_lock:
        KNOWLEDGE_DOMAINS.clear()
        KNOWLEDGE_DOMAINS.extend(new_domains)
        SYSTEM_PROMPT = build_system_prompt()
    # 重建文本缓存
    for domain in new_domains:
        abs_dp = domain.get("_abs_data_path")
        if abs_dp:
            try:
                _build_text_cache(abs_dp)
            except Exception as e:
                logger.warning("文本缓存构建失败 [%s]: %s", domain.get("name", "?"), e)
    logger.info("知识域已热加载，当前 %d 个域", len(KNOWLEDGE_DOMAINS))


class _DomainFileHandler(FileSystemEventHandler):
    """监听 knowledge/ 目录变化：domain.yaml 触发 reload，Office/PDF 文件触发缓存更新"""

    def __init__(self):
        self._timer: threading.Timer | None = None
        self._cache_timers: dict[str, threading.Timer] = {}

    def _schedule_reload(self):
        if self._timer is not None:
            self._timer.cancel()
        self._timer = threading.Timer(1.0, reload_domains)
        self._timer.daemon = True
        self._timer.start()

    def _is_domain_yaml(self, path: str) -> bool:
        return os.path.basename(path) == "domain.yaml"

    def _is_office_file(self, path: str) -> bool:
        if ".text_cache" in path:
            return False
        ext = os.path.splitext(path)[1].lower()
        return ext in _OFFICE_EXTS

    def _find_data_dir(self, path: str) -> str | None:
        """根据文件路径找到所属知识域的 data_dir"""
        with _domains_lock:
            for domain in KNOWLEDGE_DOMAINS:
                data_dir = domain.get("_abs_data_path")
                if data_dir and os.path.normpath(path).startswith(os.path.normpath(data_dir)):
                    return data_dir
        return None

    def _schedule_cache_update(self, src_path: str, deleted: bool = False):
        """防抖 1 秒后增量更新单个文件的缓存"""
        key = src_path
        if key in self._cache_timers:
            self._cache_timers[key].cancel()

        def _do_update():
            data_dir = self._find_data_dir(src_path)
            if not data_dir:
                return
            if deleted:
                cache_file = _cache_path_for(data_dir, src_path)
                if os.path.isfile(cache_file):
                    try:
                        os.remove(cache_file)
                    except Exception:
                        pass
                # 更新 meta
                meta_path = os.path.join(data_dir, ".text_cache", "_meta.json")
                rel = os.path.relpath(src_path, data_dir)
                if os.path.isfile(meta_path):
                    try:
                        with open(meta_path, "r", encoding="utf-8") as f:
                            meta = json.load(f)
                        meta.pop(rel, None)
                        with open(meta_path, "w", encoding="utf-8") as f:
                            json.dump(meta, f, ensure_ascii=False, indent=2)
                    except Exception:
                        pass
                logger.info("缓存已删除: %s", src_path)
            else:
                _update_single_cache(data_dir, src_path)
                # 更新 meta
                meta_path = os.path.join(data_dir, ".text_cache", "_meta.json")
                rel = os.path.relpath(src_path, data_dir)
                try:
                    meta = {}
                    if os.path.isfile(meta_path):
                        with open(meta_path, "r", encoding="utf-8") as f:
                            meta = json.load(f)
                    meta[rel] = os.path.getmtime(src_path)
                    with open(meta_path, "w", encoding="utf-8") as f:
                        json.dump(meta, f, ensure_ascii=False, indent=2)
                except Exception:
                    pass
                logger.info("缓存已更新: %s", src_path)
            self._cache_timers.pop(key, None)

        t = threading.Timer(1.0, _do_update)
        t.daemon = True
        t.start()
        self._cache_timers[key] = t

    def on_created(self, event):
        if event.is_directory:
            return
        if self._is_domain_yaml(event.src_path):
            self._schedule_reload()
        elif self._is_office_file(event.src_path):
            self._schedule_cache_update(event.src_path)

    def on_modified(self, event):
        if event.is_directory:
            return
        if self._is_domain_yaml(event.src_path):
            self._schedule_reload()
        elif self._is_office_file(event.src_path):
            self._schedule_cache_update(event.src_path)

    def on_deleted(self, event):
        if event.is_directory:
            return
        if self._is_domain_yaml(event.src_path):
            self._schedule_reload()
        elif self._is_office_file(event.src_path):
            self._schedule_cache_update(event.src_path, deleted=True)


_observer: Observer | None = None


def start_watcher():
    """启动 watchdog 监听 knowledge/ 目录变化"""
    global _observer
    if _observer is not None:
        return
    if not os.path.isdir(_KNOWLEDGE_DIR):
        return
    _observer = Observer()
    _observer.schedule(_DomainFileHandler(), _KNOWLEDGE_DIR, recursive=True)
    _observer.daemon = True
    _observer.start()
    logger.info("知识域文件监听已启动: %s", _KNOWLEDGE_DIR)


def build_system_prompt(session_id: str = None) -> str:
    """合并通用指令 + 各知识域 prompt 片段，生成总 system prompt"""
    if APP_MODE == "log-analyzer":
        return _build_analyzer_prompt(session_id)
    return _build_knowledge_prompt()


def _build_analyzer_prompt(session_id: str = None) -> str:
    """构建日志分析模式的 system prompt"""
    parts = [
        "你是日志分析 AI 助手，帮助开发团队快速定位服务故障根因。",
        "重要：你必须始终使用中文回答，包括工具调用过程中的所有描述和分析，禁止使用英文。",
        "",
        "能力：",
        "- 分析日志文件，提取关键错误信息",
        "- 阅读服务源代码，定位问题代码",
        "- 扫描代码发现服务间依赖关系",
        "- 识别日志截图中的内容",
        "",
        "分析流程：",
        "1. 理解用户描述的问题现象",
        "2. 如果用户提供了日志，用 read_log 分析错误和异常",
        "3. 如果用户提供了截图，直接从图片中提取关键信息",
        "4. 用 switch_service 加载相关服务代码（需要时指定版本）",
        "5. 用 search/read_file 在代码中定位问题",
        "6. 如果需要了解服务间依赖关系，用 scan_service 扫描代码中的依赖线索",
        "7. 从日志中的关键词（服务名、错误码、IP:端口）自动关联相关服务",
        "8. 信息不足时，明确告知用户需要补充什么：",
        "   - 哪个服务的日志",
        "   - 什么时间段",
        "   - 什么版本",
        "   - 相关配置信息",
        "",
        "服务管理：",
        f"- 当用户要求添加/注册服务时，先用 read_file 读取 {os.path.join(PROJECT_ROOT, 'AI_GUIDE_ANALYZER.md')} 获取完整流程",
        "- 按照指南收集信息、生成 services.yaml 配置",
        "",
    ]

    # 按业务线分组显示已注册服务
    if _analyzer_services:
        parts.append("已注册服务：")
        # 如果有业务线分组，按业务线显示
        if _analyzer_businesses:
            # 收集已分组的服务 ID
            grouped_sids = set()
            for biz_name, sids in _analyzer_businesses.items():
                parts.append(f"  【{biz_name}】")
                for sid in sids:
                    grouped_sids.add(sid)
                    svc = _analyzer_services.get(sid)
                    if svc:
                        aliases = svc.get("aliases") or []
                        alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                        parts.append(f"    • {svc.get('name', sid)} ({sid}){alias_str} [{svc.get('language', '?')}]")
                        if svc.get("description"):
                            parts.append(f"      {svc['description']}")
                    else:
                        parts.append(f"    • {sid}（未注册）")
            # 未分组的服务
            ungrouped = [sid for sid in _analyzer_services if sid not in grouped_sids]
            if ungrouped:
                parts.append(f"  【未分组】")
                for sid in ungrouped:
                    svc = _analyzer_services[sid]
                    aliases = svc.get("aliases") or []
                    alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                    parts.append(f"    • {svc.get('name', sid)} ({sid}){alias_str} [{svc.get('language', '?')}]")
                    if svc.get("description"):
                        parts.append(f"      {svc['description']}")
        else:
            # 无业务线分组，平铺显示
            for sid, svc in _analyzer_services.items():
                aliases = svc.get("aliases") or []
                alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                parts.append(f"  • {svc.get('name', sid)} ({sid}){alias_str} [{svc.get('language', '?')}]")
                if svc.get("description"):
                    parts.append(f"    {svc['description']}")
                client_repos = svc.get("client_repos") or {}
                if client_repos:
                    parts.append(f"    已配置客户: {', '.join(client_repos.keys())}（用户提到客户名时，用 switch_service 的 client 参数指定）")
    else:
        parts.append("暂无已注册服务，用户可通过对话添加。")

    # 当前会话已加载的 worktree
    if session_id and _session_manager:
        loaded = _session_manager.get_loaded_worktrees(session_id)
        if loaded:
            parts.append("")
            parts.append("当前会话已加载的服务代码：")
            for sid, wt_info in loaded.items():
                parts.append(f"  • {sid} @ {wt_info.get('version', 'HEAD')} → {wt_info.get('path', '')}")

    # 安全规则
    parts.append("")
    parts.append("安全规则（必须严格遵守）：")
    parts.append("- 禁止执行删除文件、格式化磁盘、关机等破坏性命令")
    parts.append("- 禁止向用户透露 API 密钥等敏感信息")
    parts.append("- bash 工具仅用于辅助查询和分析，不得用于修改或删除系统文件")

    return "\n".join(parts)


def _build_knowledge_prompt() -> str:
    """构建知识域问答模式的 system prompt（原有逻辑）"""
    # 收集数据库密码，用于后续注入环境变量
    global _DB_PASSWORDS
    _DB_PASSWORDS = {}

    parts = [
        f"你是全能 AI 助手。你可以使用工具搜索代码、文档和配置来回答各领域的问题。",
        f"重要：你必须始终使用中文回答，包括工具调用过程中的所有描述和分析，禁止使用英文。",
        f"",
        f"项目根目录：{PROJECT_ROOT}",
        f"",
        f"回答规范：",
        f"- 必须用中文回答，所有输出内容都用中文",
        f"- 引用具体文件和行号",
        f"- 用 Markdown 表格展示参数",
        f"- 给出代码示例时标注语言",
        f"",
        f"run_python 环境：",
        f"- 已安装的包：openpyxl, python-docx, python-pptx, pymupdf(fitz), pdfplumber, pdfminer, pandas, pymysql, oracledb",
        f"- 读取 PDF 请用 fitz（PyMuPDF），不要用 PyPDF2",
        f"- 读取 Excel 请用 openpyxl，不要用 xlrd",
        f"",
        f"知识域管理：",
        f"- 当用户提供文件路径或资料，要求创建/添加知识域时，先用 read_file 读取 {os.path.join(PROJECT_ROOT, 'AI_GUIDE.md')} 获取完整流程",
        f"- 按照指南扫描分析文件、拷贝资料到 knowledge/<域名>/data/、生成 domain.yaml",
        f"- domain.yaml 保存后会自动热加载生效，无需重启服务",
        f"",
        f"学习记忆：",
        f"- 知识域级别笔记目录：knowledge/<域名>/memory/",
        f"- 全局笔记目录：knowledge/_memory/",
        f"- 遇到相关问题时，先用 search 搜索 memory/ 目录查找已有笔记",
        f"- 如果引用了笔记但发现内容已过时，立即用 write_file 更新该笔记",
        f"",
        f"何时记录（按优先级）：",
        f"1. 用户纠正了你的错误 → 记录正确结论，防止再犯",
        f'2. 用户提供了文档中没有的隐性知识（如"这个字段虽然叫 status 但实际存的是时间戳"）',
        f"3. 经过反复搜索仍难以定位、最终通过多线索交叉验证才得出的结论 → 记录最终结果和定位路径，下次直接用",
        f"4. 可复用的查询模板、排查步骤、配置套路",
        f"",
        f"不记录：",
        f"- 单次搜索就能回答的简单问题",
        f"- 临时性数据查询结果（具体数值会变）",
        f"- 源码/文档中已明确记载的内容",
        f"",
        f"写入规则：",
        f"- 在回答用户问题的同一轮回复中，如果判断需要记录，同时调用 write_file 写入笔记（文本回答和工具调用一起返回）",
        f"- 写入前 search memory/ 检查是否已有相关笔记，有则更新同一文件而非新建",
        f"- 笔记归属明确的知识域时写到该域的 memory/，否则写到 knowledge/_memory/",
        f"- 文件名格式：YYYY-MM-DD_NNN.md（如 2026-03-08_001.md）",
        f"- 每条笔记开头标注 [学习笔记]，正文控制在 500 字以内",
        f"",
        f"你具备以下知识域，请根据用户问题自动判断所属领域并搜索对应路径：",
    ]

    for domain in KNOWLEDGE_DOMAINS:
        parts.append("")
        prompt_text = domain.get("prompt", "").strip()
        if prompt_text:
            parts.append(prompt_text)
        # 列出搜索路径供 AI 参考
        search_paths = domain.get("search_paths", [])
        if search_paths:
            parts.append(f"  搜索路径：{', '.join(search_paths)}")
        # 列出数据文件目录
        abs_data_path = domain.get("_abs_data_path")
        if abs_data_path:
            parts.append(f"  数据文件目录：{abs_data_path}")
            # 如果有 wiki 目录，显式提示搜索
            wiki_dir = os.path.join(abs_data_path, "wiki")
            if os.path.isdir(wiki_dir):
                wiki_rel = os.path.relpath(wiki_dir, PROJECT_ROOT)
                parts.append(f"  Confluence 文档目录：{wiki_rel}（包含技术文档、算法说明等，遇到相关问题务必搜索此目录）")

    # 注入数据库连接信息（密码用掩码替换）
    db_sections = []
    for domain in KNOWLEDGE_DOMAINS:
        databases = domain.get("databases")
        if not databases:
            continue
        domain_name = domain.get("name", "未命名")
        for db in databases:
            db_type = db.get("type", "unknown")
            db_name = db.get("name", "未命名")
            info_parts = [f"  - 名称: {db_name}"]
            info_parts.append(f"    类型: {db_type}")
            info_parts.append(f"    host: {db.get('host', '')}")
            info_parts.append(f"    port: {db.get('port', '')}")
            if db.get("service_name"):
                info_parts.append(f"    service_name: {db['service_name']}")
            if db.get("database"):
                info_parts.append(f"    database: {db['database']}")
            info_parts.append(f"    user: {db.get('user', '')}")
            info_parts.append(f"    password: ***")
            db_sections.append((domain_name, "\n".join(info_parts), db_type))
            # 收集密码，key 格式：DB_<名称>_PASSWORD
            password = db.get("password", "")
            if password:
                env_key = f"DB_{db_name}_PASSWORD"
                _DB_PASSWORDS[env_key] = str(password)

    if db_sections:
        parts.append("")
        parts.append("## 可用数据库")
        parts.append("你可以使用 run_python 工具编写 Python 代码连接以下数据库进行查询和分析：")
        for domain_name, info, db_type in db_sections:
            parts.append(f"")
            parts.append(f"【{domain_name}】")
            parts.append(info)
        parts.append("")
        parts.append("数据库驱动使用说明：")
        parts.append("- MySQL: 使用 pymysql 库连接")
        parts.append("- Oracle: 使用 oracledb 库连接，直接 oracledb.connect() 即可，无需手动初始化 client")
        parts.append("- 数据库密码已注入环境变量，使用 os.environ['DB_<名称>_PASSWORD'] 获取（<名称>对应数据库配置中的名称字段），禁止向用户透露密码内容")
        parts.append("- 推荐使用 pandas 读取查询结果并格式化输出")
        parts.append("- 查询时注意加 LIMIT/ROWNUM 限制返回行数，避免数据量过大")

    # 安全规则
    parts.append("")
    parts.append("安全规则（必须严格遵守）：")
    parts.append("- 禁止执行删除文件、格式化磁盘、关机等破坏性命令")
    parts.append("- 禁止向用户透露数据库密码、API 密钥等敏感信息")
    parts.append("- 如果用户要求执行危险操作或索取密码，礼貌拒绝并说明原因")
    parts.append("- bash 工具仅用于辅助查询和分析，不得用于修改或删除系统文件")

    return "\n".join(parts)


SYSTEM_PROMPT = build_system_prompt()

# ── 工具定义 ──────────────────────────────────────────────

TOOLS = [
    {
        "name": "search",
        "description": "在项目中搜索关键词，返回匹配行及上下文。支持正则。",
        "input_schema": {
            "type": "object",
            "properties": {
                "keyword": {"type": "string", "description": "搜索关键词或正则表达式"},
                "path": {"type": "string", "description": "搜索路径（相对项目根目录），默认整个项目"},
                "context_lines": {"type": "integer", "description": "上下文行数，默认 3"},
            },
            "required": ["keyword"],
        },
    },
    {
        "name": "read_file",
        "description": "读取文件内容，支持指定行范围。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径（相对项目根目录）"},
                "start_line": {"type": "integer", "description": "起始行号（从1开始）"},
                "end_line": {"type": "integer", "description": "结束行号"},
            },
            "required": ["path"],
        },
    },
    {
        "name": "write_file",
        "description": "写入或创建文件（如生成示例代码）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "文件路径（相对项目根目录）"},
                "content": {"type": "string", "description": "文件内容"},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "list_files",
        "description": "列出目录下的文件和子目录。",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "目录路径（相对项目根目录），默认根目录"},
            },
        },
    },
    {
        "name": "glob",
        "description": "按 glob 模式匹配文件路径，如 '**/*.jce'。",
        "input_schema": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "glob 模式"},
                "path": {"type": "string", "description": "起始目录（相对项目根目录），默认根目录"},
            },
            "required": ["pattern"],
        },
    },
    {
        "name": "bash",
        "description": "执行 shell 命令（如 wc、head、diff 等辅助操作）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "command": {"type": "string", "description": "要执行的命令"},
            },
            "required": ["command"],
        },
    },
    {
        "name": "web_fetch",
        "description": "抓取网页内容（如在线文档、API 说明）。",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "网页 URL"},
            },
            "required": ["url"],
        },
    },
    {
        "name": "run_python",
        "description": "执行 Python 代码，可用于数据库查询、数据分析、数据处理等。",
        "input_schema": {
            "type": "object",
            "properties": {
                "code": {"type": "string", "description": "要执行的 Python 代码"},
            },
            "required": ["code"],
        },
    },
]

# ── 日志分析模式专用工具 ──────────────────────────────────

_ANALYZER_TOOLS = [
    {
        "name": "read_log",
        "description": "读取并过滤日志文件，支持按级别、关键词、时间范围筛选。",
        "input_schema": {
            "type": "object",
            "properties": {
                "file": {"type": "string", "description": "文件路径（相对 session uploads 目录）"},
                "level": {"type": "string", "description": "过滤级别：ERROR/WARN/INFO（可选）"},
                "keyword": {"type": "string", "description": "关键词过滤，支持正则（可选）"},
                "time_start": {"type": "string", "description": "起始时间，如 2026-03-20 15:00:00（可选）"},
                "time_end": {"type": "string", "description": "结束时间（可选）"},
                "context_lines": {"type": "integer", "description": "上下文行数，默认 3"},
                "tail": {"type": "integer", "description": "只看最后 N 行（可选）"},
            },
            "required": ["file"],
        },
    },
    {
        "name": "scan_service",
        "description": "扫描服务代码，发现依赖关系线索（配置文件中的服务引用、RPC接口定义、构建文件依赖等）。需要先用 switch_service 加载服务代码。",
        "input_schema": {
            "type": "object",
            "properties": {
                "service": {"type": "string", "description": "服务 ID"},
            },
            "required": ["service"],
        },
    },
    {
        "name": "switch_service",
        "description": "加载服务代码到当前会话，后续可用 search/read_file 查看代码。如果用户提到了客户名，传 client 参数自动匹配该客户的仓库。",
        "input_schema": {
            "type": "object",
            "properties": {
                "service": {"type": "string", "description": "服务 ID"},
                "version": {"type": "string", "description": "版本：分支名/tag/commit hash（可选，默认最新代码）"},
                "client": {"type": "string", "description": "客户名（可选），匹配该客户专属的代码仓库"},
            },
            "required": ["service"],
        },
    },
    {
        "name": "list_services",
        "description": "列出所有已注册的服务，按业务线分组显示。",
        "input_schema": {
            "type": "object",
            "properties": {},
        },
    },
]

if APP_MODE == "log-analyzer":
    TOOLS += _ANALYZER_TOOLS
    # 移除日志分析模式不需要的工具
    TOOLS = [t for t in TOOLS if t["name"] != "write_file"]
    # 调整 search 工具描述，限定搜索范围为已加载的服务代码
    for t in TOOLS:
        if t["name"] == "search":
            t["description"] = "在已加载的服务代码中搜索关键词，返回匹配行及上下文。需要先用 switch_service 加载服务代码。支持正则。"
            break
    # 扩展必填参数映射
    _REQUIRED_FIELDS_EXTRA = {
        "read_log": ["file"],
        "scan_service": ["service"],
        "switch_service": ["service"],
    }

# ── 工具执行 ──────────────────────────────────────────────


def _safe_path(rel: str, allowed_paths: list[str] = None) -> str:
    """将路径转为安全的绝对路径。

    Args:
        rel: 相对或绝对路径
        allowed_paths: 允许访问的目录列表（log-analyzer 模式），为 None 时限制在 PROJECT_ROOT
    """
    if allowed_paths:
        # log-analyzer 模式：限制在 session 的 worktree + uploads 内
        if not rel or rel in (".", "/", ""):
            # 空路径：优先返回第一个 worktree 路径（跳过 uploads）
            return allowed_paths[-1] if len(allowed_paths) > 1 else allowed_paths[0]
        abs_path = os.path.normpath(rel) if os.path.isabs(rel) else os.path.normpath(os.path.join(allowed_paths[-1], rel))
        if any(abs_path.startswith(p) for p in allowed_paths):
            return abs_path
        # 不在允许范围内，返回最后一个 worktree 路径
        return allowed_paths[-1] if len(allowed_paths) > 1 else allowed_paths[0]

    abs_path = os.path.normpath(os.path.join(PROJECT_ROOT, rel))
    if not abs_path.startswith(PROJECT_ROOT):
        return PROJECT_ROOT
    return abs_path


def _read_office_file(fpath: str) -> str:
    """根据扩展名解析 Office 文件，返回纯文本内容"""
    ext = os.path.splitext(fpath)[1].lower()

    if ext in ('.xlsx', '.xls'):
        import openpyxl
        wb = openpyxl.load_workbook(fpath, read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"=== Sheet: {sheet} ===")
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join(str(c) if c is not None else "" for c in row))
        wb.close()
        return "\n".join(parts)

    if ext == '.docx':
        import docx
        doc = docx.Document(fpath)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(fpath)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if texts:
                parts.append(f"=== Slide {i} ===")
                parts.append("\n".join(texts))
        return "\n".join(parts)

    if ext == '.pdf':
        try:
            import fitz  # PyMuPDF，速度远快于 pdfplumber
            parts = []
            doc = fitz.open(fpath)
            for i, page in enumerate(doc, 1):
                text = page.get_text()
                if text and text.strip():
                    parts.append(f"=== Page {i} ===")
                    parts.append(text)
            doc.close()
            return "\n".join(parts)
        except ImportError:
            import pdfplumber
            parts = []
            with pdfplumber.open(fpath) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        parts.append(f"=== Page {i} ===")
                        parts.append(text)
            return "\n".join(parts)

    return ""


# ── Office/PDF 文本缓存 ──────────────────────────────────

def _cache_path_for(data_dir: str, src_path: str) -> str:
    """返回源文件对应的缓存 txt 路径"""
    rel = os.path.relpath(src_path, data_dir)
    return os.path.join(data_dir, ".text_cache", rel + ".txt")


def _update_single_cache(data_dir: str, src_path: str) -> None:
    """对单个 Office/PDF 文件生成或更新缓存，使用临时文件+rename 保证原子写入"""
    cache_file = _cache_path_for(data_dir, src_path)
    try:
        text = _read_office_file(src_path)
        if text:
            os.makedirs(os.path.dirname(cache_file), exist_ok=True)
            tmp_file = cache_file + ".tmp"
            with open(tmp_file, "w", encoding="utf-8") as f:
                f.write(text)
            os.replace(tmp_file, cache_file)
    except Exception as e:
        # 清理可能残留的临时文件
        tmp_file = cache_file + ".tmp"
        if os.path.isfile(tmp_file):
            try:
                os.remove(tmp_file)
            except Exception:
                pass
        logger.warning("缓存生成失败 %s: %s", src_path, e)


def _save_meta(meta_path: str, meta: dict) -> None:
    """原子写入 meta 文件"""
    tmp = meta_path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    os.replace(tmp, meta_path)


def _build_text_cache(data_dir: str) -> None:
    """扫描 data_dir 下所有 Office/PDF 文件，增量生成纯文本缓存"""
    if not os.path.isdir(data_dir):
        return

    cache_dir = os.path.join(data_dir, ".text_cache")
    meta_path = os.path.join(cache_dir, "_meta.json")

    # 读取已有元数据
    meta: dict[str, float] = {}
    if os.path.isfile(meta_path):
        try:
            with open(meta_path, "r", encoding="utf-8") as f:
                meta = json.load(f)
        except Exception:
            meta = {}

    # 扫描所有 Office/PDF 文件
    current_files: dict[str, float] = {}
    for root, dirs, files in os.walk(data_dir):
        # 跳过 .text_cache 目录本身
        dirs[:] = [d for d in dirs if d != ".text_cache"]
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if ext in _OFFICE_EXTS:
                fpath = os.path.join(root, fname)
                rel = os.path.relpath(fpath, data_dir)
                current_files[rel] = os.path.getmtime(fpath)

    # 增量更新：新增、mtime 变化、或缓存文件不存在的
    need_update = []
    for rel, mtime in current_files.items():
        cache_file = os.path.join(cache_dir, rel + ".txt")
        if rel not in meta or meta[rel] != mtime or not os.path.isfile(cache_file):
            need_update.append((rel, mtime))

    os.makedirs(cache_dir, exist_ok=True)

    if need_update:
        total = len(need_update)
        logger.info("文本缓存：需更新 %d 个文件 (%s)", total, data_dir)

        from concurrent.futures import ProcessPoolExecutor, as_completed
        import multiprocessing
        workers = min(len(need_update), max(2, multiprocessing.cpu_count()))
        done_count = 0

        with ProcessPoolExecutor(max_workers=workers) as pool:
            # 子进程只做解析，返回 (rel, mtime, text)
            futures = {}
            for rel, mtime in need_update:
                src_path = os.path.join(data_dir, rel)
                futures[pool.submit(_read_office_file, src_path)] = (rel, mtime)

            for future in as_completed(futures):
                rel, mtime = futures[future]
                try:
                    text = future.result()
                    if text:
                        cache_file = os.path.join(cache_dir, rel + ".txt")
                        os.makedirs(os.path.dirname(cache_file), exist_ok=True)
                        tmp_file = cache_file + ".tmp"
                        with open(tmp_file, "w", encoding="utf-8") as f:
                            f.write(text)
                        os.replace(tmp_file, cache_file)
                    meta[rel] = mtime
                    done_count += 1
                    logger.info("文本缓存 [%d/%d]: %s", done_count, total, rel)
                    _save_meta(meta_path, meta)
                except Exception as e:
                    logger.warning("文本缓存处理失败 %s: %s", rel, e)

    # 清理：源文件已删除的缓存
    cleaned = False
    for rel in list(meta.keys()):
        if rel not in current_files:
            cache_file = os.path.join(cache_dir, rel + ".txt")
            if os.path.isfile(cache_file):
                try:
                    os.remove(cache_file)
                except Exception:
                    pass
            del meta[rel]
            cleaned = True

    if cleaned:
        _save_meta(meta_path, meta)

    if current_files:
        logger.info("文本缓存已就绪: %s (%d 个文件)", data_dir, len(current_files))


def _restore_cache_path(grep_line: str, data_dir: str) -> str:
    """将 grep 输出中的缓存路径还原为原始 Office/PDF 文件路径
    例：/path/data/.text_cache/sub/报告.pdf.txt:10:内容 → knowledge/xxx/data/sub/报告.pdf:10:内容
    """
    cache_dir = os.path.join(data_dir, ".text_cache")
    # grep 输出格式：文件路径:行号:内容 或 文件路径-行号-内容（上下文行）
    if cache_dir in grep_line:
        # 替换 .text_cache/ 路径为 data_dir 路径，去掉 .txt 后缀
        restored = grep_line.replace(cache_dir + os.sep, data_dir + os.sep)
        # 处理 .pdf.txt:行号 → .pdf:行号
        restored = re.sub(r'\.txt([:|-])', r'\1', restored, count=1)
        # 转为相对路径显示
        abs_prefix = data_dir + os.sep
        if restored.startswith(abs_prefix):
            rel_prefix = os.path.relpath(data_dir, PROJECT_ROOT) + os.sep
            restored = rel_prefix + restored[len(abs_prefix):]
        return restored
    return grep_line


def _find_cache_file(fpath: str) -> str | None:
    """查找 Office/PDF 文件对应的缓存 txt，存在则返回路径，否则 None"""
    with _domains_lock:
        for domain in KNOWLEDGE_DOMAINS:
            data_dir = domain.get("_abs_data_path")
            if data_dir and os.path.normpath(fpath).startswith(os.path.normpath(data_dir)):
                cache_file = _cache_path_for(data_dir, fpath)
                if os.path.isfile(cache_file):
                    return cache_file
                return None
    return None


# ── 启动时构建文本缓存 ────────────────────────────────────
def init_text_cache():
    """初始化所有知识域的文本缓存，应在 logging 配置完成后调用"""
    for domain in KNOWLEDGE_DOMAINS:
        abs_dp = domain.get("_abs_data_path")
        if abs_dp:
            try:
                _build_text_cache(abs_dp)
            except Exception as e:
                logger.warning("文本缓存构建失败 [%s]: %s", domain.get("name", "?"), e)


# ── 危险命令黑名单 ────────────────────────────────────────

_DANGEROUS_PATTERNS = [
    (re.compile(r'\brm\s+-\S*r', re.IGNORECASE), '禁止递归删除'),
    (re.compile(r'\brm\s+-\S*f', re.IGNORECASE), '禁止强制删除'),
    (re.compile(r'\brm\s+.*\s/\s*$', re.IGNORECASE), '禁止删除根目录'),
    (re.compile(r'\bmkfs\b', re.IGNORECASE), '禁止格式化磁盘'),
    (re.compile(r'\bdd\s+', re.IGNORECASE), '禁止 dd 命令'),
    (re.compile(r'(?<!\d)>\s*/dev/(?!null)', re.IGNORECASE), '禁止写设备文件'),
    (re.compile(r'\bchmod\s+777\b', re.IGNORECASE), '禁止设置 777 权限'),
    (re.compile(r':\(\)\s*\{'), '禁止 fork bomb'),
    (re.compile(r'\b(shutdown|reboot|poweroff|halt)\b', re.IGNORECASE), '禁止关机重启'),
    (re.compile(r'(curl|wget)\s.*\|\s*(ba)?sh', re.IGNORECASE), '禁止远程代码执行'),
    (re.compile(r'curl\s.*-d\s+@', re.IGNORECASE), '禁止通过 curl 外泄数据'),
    (re.compile(r'curl\s.*--data\s+@', re.IGNORECASE), '禁止通过 curl 外泄数据'),
    (re.compile(r'base64\s.*\|\s*(ba)?sh', re.IGNORECASE), '禁止 base64 解码执行'),
    (re.compile(r'\beval\b', re.IGNORECASE), '禁止 eval 执行'),
    (re.compile(r'\$\(.*\)\s*\|\s*(ba)?sh', re.IGNORECASE), '禁止命令替换执行'),
]

# run_python 危险代码模式
_DANGEROUS_PYTHON_PATTERNS = [
    (re.compile(r'\bos\.system\b'), '禁止 os.system 调用'),
    (re.compile(r'\bos\.popen\b'), '禁止 os.popen 调用'),
    (re.compile(r'\bos\.exec'), '禁止 os.exec* 调用'),
    (re.compile(r'\bos\.remove\b'), '禁止 os.remove 调用'),
    (re.compile(r'\bos\.unlink\b'), '禁止 os.unlink 调用'),
    (re.compile(r'\bos\.rmdir\b'), '禁止 os.rmdir 调用'),
    (re.compile(r'\bshutil\.rmtree\b'), '禁止 shutil.rmtree 调用'),
    (re.compile(r'\bsubprocess\b'), '禁止 subprocess 调用'),
    (re.compile(r'\b__import__\b'), '禁止动态导入'),
]


def _check_dangerous_command(cmd: str) -> str | None:
    """检查命令是否包含危险模式，返回拒绝原因或 None"""
    # 去除反斜杠转义绕过
    normalized = cmd.replace('\\', '')
    for pattern, reason in _DANGEROUS_PATTERNS:
        if pattern.search(normalized):
            return reason
    return None


def _check_dangerous_python(code: str) -> str | None:
    """检查 Python 代码是否包含危险模式，返回拒绝原因或 None"""
    for pattern, reason in _DANGEROUS_PYTHON_PATTERNS:
        if pattern.search(code):
            return reason
    return None


def _sanitize_output(text: str) -> str:
    """将工具输出中的敏感信息替换为掩码"""
    # 替换数据库密码
    for env_key, env_val in _DB_PASSWORDS.items():
        if env_val and env_val in text:
            text = text.replace(env_val, '***')
    # 替换 API Key
    if API_KEY and API_KEY in text:
        text = text.replace(API_KEY, '***')
    return text


# 工具必填参数映射
_REQUIRED_FIELDS = {
    "search": ["keyword"],
    "read_file": ["path"],
    "write_file": ["path", "content"],
    "glob": ["pattern"],
    "bash": ["command"],
    "web_fetch": ["url"],
    "run_python": ["code"],
}

if APP_MODE == "log-analyzer":
    _REQUIRED_FIELDS.update(_REQUIRED_FIELDS_EXTRA)


def _check_tool_params(name: str, inp: dict) -> str | None:
    """校验工具必填参数，返回错误消息或 None"""
    missing = [f for f in _REQUIRED_FIELDS.get(name, []) if f not in inp or not inp[f]]
    if missing:
        return f"参数错误：缺少必填参数 {', '.join(missing)}，请重新调用并提供完整参数"
    return None


def exec_tool(name: str, inp: dict) -> str:
    """执行工具，返回结果字符串"""
    logger.info("执行工具: %s, 参数: %s", name, json.dumps(inp, ensure_ascii=False)[:500])

    # 校验必填参数
    param_err = _check_tool_params(name, inp)
    if param_err:
        return param_err

    t0 = time.time()
    # log-analyzer 模式下，限制文件操作在当前 session 的允许路径内
    _allowed = None
    if APP_MODE == "log-analyzer" and _session_manager:
        _sid = inp.get("_session_id", "")
        if _sid:
            _allowed = _session_manager.get_allowed_paths(_sid)
    try:
        if name == "search":
            keyword = inp["keyword"]
            path = _safe_path(inp.get("path", ""), _allowed)
            ctx = str(inp.get("context_lines", 3))
            # log-analyzer 模式搜索服务代码，需要更多文件类型
            if APP_MODE == "log-analyzer":
                include_args = [
                    "--include=*.h", "--include=*.cpp", "--include=*.c", "--include=*.hpp",
                    "--include=*.java", "--include=*.py", "--include=*.go", "--include=*.rs",
                    "--include=*.js", "--include=*.ts", "--include=*.cs",
                    "--include=*.proto", "--include=*.thrift", "--include=*.jce",
                    "--include=*.xml", "--include=*.yaml", "--include=*.yml",
                    "--include=*.json", "--include=*.ini", "--include=*.conf",
                    "--include=*.properties", "--include=*.toml",
                    "--include=*.md", "--include=*.txt", "--include=*.sh",
                    "--include=*.cmake", "--include=CMakeLists.txt",
                    "--include=*.gradle", "--include=*.sln", "--include=*.csproj",
                    "--include=Makefile", "--include=*.mk",
                ]
            else:
                include_args = [
                    "--include=*.jce",
                    "--include=*.h", "--include=*.cpp", "--include=*.md",
                    "--include=*.conf", "--include=*.xml", "--include=*.yaml",
                    "--include=*.yml", "--include=*.txt", "--include=*.sh",
                ]
            result = subprocess.run(
                ["grep", "-E", "-r", "-n", f"-C{ctx}",
                 "--exclude-dir=logs", "--exclude-dir=shares",
                 "--exclude-dir=.venv", "--exclude-dir=__pycache__", "--exclude-dir=.git",
                 "--exclude-dir=.text_cache",
                 *include_args,
                 keyword, path],
                capture_output=True, text=True, timeout=30,
            )
            output = result.stdout or result.stderr or "无匹配结果"

            # 搜索 Office/PDF 文本缓存
            cache_matches = []
            for cache_dir_path in glob_mod.glob(os.path.join(path, "**", ".text_cache"), recursive=True):
                if not os.path.isdir(cache_dir_path):
                    continue
                cache_result = subprocess.run(
                    ["grep", "-E", "-r", "-n", f"-C{ctx}",
                     "--include=*.txt",
                     keyword, cache_dir_path],
                    capture_output=True, text=True, timeout=30,
                )
                if cache_result.stdout:
                    # 还原路径：.text_cache/xxx.pdf.txt:行号 → xxx.pdf:行号
                    data_dir = os.path.dirname(cache_dir_path)
                    for line in cache_result.stdout.splitlines():
                        restored = _restore_cache_path(line, data_dir)
                        cache_matches.append(restored)
            if cache_matches:
                if output == "无匹配结果":
                    output = "\n".join(cache_matches)
                else:
                    output += "\n" + "\n".join(cache_matches)

            # 子目录搜索无结果时，自动扩大到整个 knowledge/ 目录重搜（仅 knowledge 模式）
            if APP_MODE == "knowledge":
                knowledge_dir = os.path.join(PROJECT_ROOT, "knowledge")
                if output == "无匹配结果" and path != PROJECT_ROOT and path != knowledge_dir:
                    fallback = subprocess.run(
                        ["grep", "-E", "-r", "-n", "-C1",
                         "--exclude-dir=logs", "--exclude-dir=shares",
                         "--exclude-dir=.venv", "--exclude-dir=__pycache__", "--exclude-dir=.git",
                         "--include=*.md", "--include=*.txt", "--include=*.yaml",
                         keyword, knowledge_dir],
                        capture_output=True, text=True, timeout=30,
                    )
                    if fallback.stdout:
                        output = f"在 {os.path.relpath(path, PROJECT_ROOT)} 中未找到，已自动扩大搜索范围：\n{fallback.stdout}"

        elif name == "read_file":
            fpath = _safe_path(inp["path"], _allowed)
            ext = os.path.splitext(fpath)[1].lower()
            if ext in _OFFICE_EXTS:
                # 优先读缓存
                cache_file = _find_cache_file(fpath)
                if cache_file:
                    with open(cache_file, "r", encoding="utf-8", errors="replace") as f:
                        lines = f.readlines()
                    start = max(1, inp.get("start_line", 1))
                    end = inp.get("end_line", len(lines))
                    selected = lines[start - 1:end]
                    output = "".join(f"{start + i}: {l}" for i, l in enumerate(selected))
                else:
                    output = _read_office_file(fpath) or "(空文件)"
            else:
                with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                    lines = f.readlines()
                start = max(1, inp.get("start_line", 1))
                end = inp.get("end_line", len(lines))
                selected = lines[start - 1:end]
                output = "".join(f"{start + i}: {l}" for i, l in enumerate(selected))

        elif name == "write_file":
            fpath = _safe_path(inp["path"])
            os.makedirs(os.path.dirname(fpath), exist_ok=True)
            with open(fpath, "w", encoding="utf-8") as f:
                f.write(inp["content"])
            output = f"已写入 {fpath}"

        elif name == "list_files":
            dpath = _safe_path(inp.get("path", ""), _allowed)
            entries = sorted(os.listdir(dpath))
            output = "\n".join(entries) if entries else "空目录"

        elif name == "glob":
            base = _safe_path(inp.get("path", ""), _allowed)
            pattern = inp["pattern"]
            matches = sorted(glob_mod.glob(os.path.join(base, pattern), recursive=True))
            rel = [os.path.relpath(m, PROJECT_ROOT) for m in matches]
            output = "\n".join(rel) if rel else "无匹配文件"

        elif name == "bash":
            cmd = inp["command"]
            danger = _check_dangerous_command(cmd)
            if danger:
                logger.warning("bash 命令被安全拦截: %s, 原因: %s", cmd, danger)
                output = "该命令因安全策略被拒绝执行，请使用其他方式完成任务"
            else:
                bash_env = os.environ.copy()
                venv_bin = os.path.join(PROJECT_ROOT, ".venv", "bin")
                if os.path.isdir(venv_bin):
                    bash_env["PATH"] = venv_bin + ":" + bash_env.get("PATH", "")
                bash_cwd = _allowed[0] if _allowed else PROJECT_ROOT
                result = subprocess.run(
                    cmd, shell=True, capture_output=True, text=True,
                    timeout=30, cwd=bash_cwd, env=bash_env,
                )
                output = (result.stdout + result.stderr).strip() or "(无输出)"

        elif name == "web_fetch":
            req = urllib.request.Request(inp["url"], headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=15) as resp:
                output = resp.read().decode("utf-8", errors="replace")

        elif name == "run_python":
            import tempfile
            code = inp["code"]
            # 安全检查
            danger = _check_dangerous_python(code)
            if danger:
                logger.warning("Python 代码被安全拦截: %s, 原因: %s", code[:200], danger)
                output = "该代码因安全策略被拒绝执行，请使用其他方式完成任务"
            else:
                # 如果配置了 Oracle Client 路径，自动注入初始化代码
                if ORACLE_CLIENT_PATH and "oracledb" in code and "init_oracle_client" not in code:
                    code = (
                        "import oracledb\n"
                        f"oracledb.init_oracle_client(lib_dir={ORACLE_CLIENT_PATH!r})\n"
                        + code
                    )
                with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, encoding="utf-8") as f:
                    f.write(code)
                    tmp_path = f.name
                try:
                    timeout = CONFIG["tools"].get("python_timeout", 300)
                    env = os.environ.copy()
                    if ORACLE_CLIENT_PATH:
                        env["LD_LIBRARY_PATH"] = ORACLE_CLIENT_PATH + ":" + env.get("LD_LIBRARY_PATH", "")
                    # 注入数据库密码为环境变量
                    for env_key, env_val in _DB_PASSWORDS.items():
                        env[env_key] = env_val
                    result = subprocess.run(
                        ["python3", tmp_path],
                        capture_output=True, text=True,
                        timeout=timeout, cwd=PROJECT_ROOT, env=env,
                    )
                    output = (result.stdout + result.stderr).strip() or "(无输出)"
                except subprocess.TimeoutExpired:
                    output = f"执行超时（{timeout}秒），代码已终止"
                finally:
                    os.unlink(tmp_path)

        # ── 日志分析模式专用工具 ──────────────────
        elif name == "read_log" and APP_MODE == "log-analyzer":
            session_id = inp.get("_session_id", "")
            file_rel = inp["file"]
            if session_id and _session_manager:
                uploads_dir = _session_manager.get_uploads_path(session_id)
                filepath = os.path.join(uploads_dir, file_rel)
            else:
                filepath = file_rel
            output = read_log_filtered(
                filepath,
                level=inp.get("level"),
                keyword=inp.get("keyword"),
                time_start=inp.get("time_start"),
                time_end=inp.get("time_end"),
                context_lines=inp.get("context_lines", 3),
                tail=inp.get("tail"),
            )

        elif name == "scan_service" and APP_MODE == "log-analyzer":
            service_id = inp["service"]
            # 支持别名查找
            resolved_id, svc = _find_service(service_id)
            if resolved_id:
                service_id = resolved_id
            if not svc:
                output = f"未找到服务: {inp['service']}，请用 list_services 查看已注册服务"
            else:
                session_id = inp.get("_session_id", "")
                if not session_id or not _session_manager:
                    output = "无法扫描服务代码：缺少 session 信息"
                else:
                    loaded_wts = _session_manager.get_loaded_worktrees(session_id)
                    wt_info = loaded_wts.get(service_id)
                    if not wt_info:
                        output = f"服务 {service_id} 的代码尚未加载，请先用 switch_service 加载"
                    else:
                        code_path = wt_info.get("path", "")
                        sub_path = wt_info.get("sub_path")
                        scan_path = os.path.join(code_path, sub_path) if sub_path else code_path
                        output = scan_service_deps(scan_path, _analyzer_services)

        elif name == "switch_service" and APP_MODE == "log-analyzer":
            service_id = inp["service"]
            version = inp.get("version")
            client = inp.get("client")
            session_id = inp.get("_session_id", "")
            # 支持别名查找
            resolved_id, svc = _find_service(service_id)
            if resolved_id:
                service_id = resolved_id
            if not svc:
                output = f"未找到服务: {inp['service']}，请用 list_services 查看已注册服务"
            elif not session_id or not _session_manager:
                output = "无法加载服务代码：缺少 session 信息"
            else:
                repo = svc["repo"]
                sub_path = svc.get("sub_path")
                client_repos = svc.get("client_repos") or {}
                code_path = _session_manager.setup_code(
                    session_id, service_id, repo, version, sub_path,
                    client=client, client_repos=client_repos,
                )
                # 获取实际的 sub_path（可能由 URL 解析或 client_repos 覆盖）
                loaded_wt = _session_manager.get_loaded_worktrees(session_id).get(service_id, {})
                actual_sub = loaded_wt.get("sub_path")
                output = f"已加载 {svc.get('name', service_id)} 代码到: {code_path}\n"
                output += f"语言: {svc.get('language', '未知')}\n"
                output += f"版本: {version or '最新'}\n"
                if actual_sub:
                    output += f"该服务主代码目录: {code_path}/{actual_sub}\n"
                    output += f"仓库根目录: {code_path}（可搜索整个仓库，包括公共库等其他目录）\n"
                if client:
                    used_repo = client_repos.get(client, repo)
                    if isinstance(used_repo, dict):
                        used_repo = used_repo.get("repo", repo)
                    output += f"客户: {client}（仓库: {used_repo}）\n"
                if client_repos:
                    output += f"已配置客户: {', '.join(client_repos.keys())}\n"
                output += f"你现在可以用 search 和 read_file 工具查看代码，路径前缀: {code_path}"

        elif name == "list_services" and APP_MODE == "log-analyzer":
            if not _analyzer_services:
                output = "暂无已注册服务，请先配置 services.yaml"
            else:
                lines = []
                if _analyzer_businesses:
                    grouped_sids = set()
                    for biz_name, sids in _analyzer_businesses.items():
                        lines.append(f"【{biz_name}】")
                        for sid in sids:
                            grouped_sids.add(sid)
                            svc = _analyzer_services.get(sid)
                            if svc:
                                aliases = svc.get("aliases") or []
                                alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                                lines.append(f"  • {svc.get('name', sid)} ({sid}){alias_str}")
                                lines.append(f"    语言: {svc.get('language', '未知')}")
                                lines.append(f"    {svc.get('description', '')}")
                                lines.append(f"    默认仓库: {svc.get('repo', '')}")
                                if svc.get("sub_path"):
                                    lines.append(f"    子路径: {svc['sub_path']}")
                                client_repos = svc.get("client_repos") or {}
                                if client_repos:
                                    lines.append(f"    已配置客户: {', '.join(client_repos.keys())}")
                            else:
                                lines.append(f"  • {sid}（未注册）")
                            lines.append("")
                    ungrouped = [sid for sid in _analyzer_services if sid not in grouped_sids]
                    if ungrouped:
                        lines.append("【未分组】")
                        for sid in ungrouped:
                            svc = _analyzer_services[sid]
                            aliases = svc.get("aliases") or []
                            alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                            lines.append(f"  • {svc.get('name', sid)} ({sid}){alias_str}")
                            lines.append(f"    语言: {svc.get('language', '未知')}")
                            lines.append(f"    {svc.get('description', '')}")
                            lines.append(f"    默认仓库: {svc.get('repo', '')}")
                            if svc.get("sub_path"):
                                lines.append(f"    子路径: {svc['sub_path']}")
                            client_repos = svc.get("client_repos") or {}
                            if client_repos:
                                lines.append(f"    已配置客户: {', '.join(client_repos.keys())}")
                            lines.append("")
                else:
                    for sid, svc in _analyzer_services.items():
                        aliases = svc.get("aliases") or []
                        alias_str = f" (别名: {', '.join(aliases)})" if aliases else ""
                        lines.append(f"• {svc.get('name', sid)} ({sid}){alias_str}")
                        lines.append(f"  语言: {svc.get('language', '未知')}")
                        lines.append(f"  {svc.get('description', '')}")
                        lines.append(f"  默认仓库: {svc.get('repo', '')}")
                        if svc.get("sub_path"):
                            lines.append(f"  子路径: {svc['sub_path']}")
                        client_repos = svc.get("client_repos") or {}
                        if client_repos:
                            lines.append(f"  已配置客户: {', '.join(client_repos.keys())}")
                        lines.append("")
                output = "\n".join(lines)

        else:
            output = f"未知工具: {name}"

    except Exception as e:
        logger.error("工具执行错误: %s, %s", name, e, exc_info=True)
        output = f"工具执行错误: {e}"

    elapsed = time.time() - t0
    if len(output) > MAX_OUTPUT_LEN:
        output = output[:MAX_OUTPUT_LEN] + f"\n... (截断，共 {len(output)} 字符)"
    # 输出脱敏：替换敏感信息
    output = _sanitize_output(output)
    logger.info("工具完成: %s, 耗时: %.2fs, 结果长度: %d", name, elapsed, len(output))
    return output


# ── Agent 流式调用 ────────────────────────────────────────


_MAX_EMPTY_RETRIES = 3  # 连续空参数工具调用的最大容忍次数


def _tools_to_openai() -> list[dict]:
    """将 Anthropic 格式的工具定义转换为 OpenAI 格式"""
    return [{
        "type": "function",
        "function": {
            "name": tool["name"],
            "description": tool["description"],
            "parameters": tool["input_schema"],
        },
    } for tool in TOOLS]


def _truncate_tool_content(content: str) -> tuple[str, int]:
    """截断单个工具结果，返回 (截断后内容, 节省字符数)"""
    orig_len = len(content)
    if orig_len <= _COMPACT_MAX_LEN:
        return content, 0
    truncated = content[:_COMPACT_MAX_LEN] + f"\n...(已截断，原始 {orig_len} 字符，如需重新查看请再次调用工具)"
    return truncated, orig_len - len(truncated)


def _truncate_tool_results_anthropic(messages: list, iteration: int):
    """截断 Anthropic 格式消息中旧轮次的工具结果（原地修改）"""
    # 从后往前找包含 tool_result 的 user 消息，标记轮次
    tool_round_indices = []  # [(user_msg_index, prev_assistant_index), ...]
    for i in range(len(messages) - 1, -1, -1):
        msg = messages[i]
        if msg["role"] != "user" or not isinstance(msg.get("content"), list):
            continue
        has_tool_result = any(
            isinstance(b, dict) and b.get("type") == "tool_result"
            for b in msg["content"]
        )
        if not has_tool_result:
            continue
        # 找前一个 assistant 消息
        prev_asst = None
        for j in range(i - 1, -1, -1):
            if messages[j]["role"] == "assistant":
                prev_asst = j
                break
        tool_round_indices.append((i, prev_asst))

    if len(tool_round_indices) <= _COMPACT_KEEP_RECENT:
        return

    # 需要截断的轮次（跳过最近 keep_recent 轮）
    rounds_to_truncate = tool_round_indices[_COMPACT_KEEP_RECENT:]
    truncated_count = 0
    saved_chars = 0

    for user_idx, asst_idx in rounds_to_truncate:
        # 从 assistant 消息建立 tool_use_id → name 映射
        id_to_name = {}
        if asst_idx is not None:
            asst_content = messages[asst_idx].get("content", [])
            if isinstance(asst_content, list):
                for b in asst_content:
                    if hasattr(b, "type") and b.type == "tool_use":
                        id_to_name[b.id] = b.name
                    elif isinstance(b, dict) and b.get("type") == "tool_use":
                        id_to_name[b.get("id")] = b.get("name")

        # 截断 tool_result
        for block in messages[user_idx]["content"]:
            if not isinstance(block, dict) or block.get("type") != "tool_result":
                continue
            tool_name = id_to_name.get(block.get("tool_use_id"), "")
            if tool_name in _COMPACT_SKIP_TOOLS:
                continue
            content = block.get("content", "")
            if isinstance(content, str) and len(content) > _COMPACT_MAX_LEN:
                block["content"], saved = _truncate_tool_content(content)
                saved_chars += saved
                truncated_count += 1

    if truncated_count > 0:
        logger.info("上下文截断: 第 %d 轮，截断 %d 个工具结果，节省约 %d 字符",
                     iteration, truncated_count, saved_chars)


def _truncate_tool_results_openai(oai_messages: list, iteration: int):
    """截断 OpenAI 格式消息中旧轮次的工具结果（原地修改）"""
    # 从后往前找 assistant+tool_calls 消息，标记轮次
    tool_round_indices = []  # [assistant_msg_index, ...]
    for i in range(len(oai_messages) - 1, -1, -1):
        msg = oai_messages[i]
        if msg.get("role") == "assistant" and msg.get("tool_calls"):
            tool_round_indices.append(i)

    if len(tool_round_indices) <= _COMPACT_KEEP_RECENT:
        return

    # 需要截断的轮次（跳过最近 keep_recent 轮）
    rounds_to_truncate = tool_round_indices[_COMPACT_KEEP_RECENT:]

    # 建立所有需截断轮次的 tool_call_id → name 映射
    ids_to_truncate = {}  # tool_call_id → tool_name
    for asst_idx in rounds_to_truncate:
        for tc in oai_messages[asst_idx].get("tool_calls", []):
            fn = tc.get("function", {})
            ids_to_truncate[tc["id"]] = fn.get("name", "")

    # 遍历所有 tool 消息，截断匹配的
    truncated_count = 0
    saved_chars = 0
    for msg in oai_messages:
        if msg.get("role") != "tool":
            continue
        tc_id = msg.get("tool_call_id")
        if tc_id not in ids_to_truncate:
            continue
        tool_name = ids_to_truncate[tc_id]
        if tool_name in _COMPACT_SKIP_TOOLS:
            continue
        content = msg.get("content", "")
        if isinstance(content, str) and len(content) > _COMPACT_MAX_LEN:
            msg["content"], saved = _truncate_tool_content(content)
            saved_chars += saved
            truncated_count += 1

    if truncated_count > 0:
        logger.info("上下文截断: 第 %d 轮，截断 %d 个工具结果，节省约 %d 字符",
                     iteration, truncated_count, saved_chars)


def _run_anthropic_stream(messages: list, session_id: str = None):
    """Anthropic SDK 流式调用"""
    import anthropic
    client = anthropic.Anthropic(base_url=BASE_URL, api_key=API_KEY)
    empty_retries = 0
    sys_prompt = build_system_prompt(session_id) if APP_MODE == "log-analyzer" else SYSTEM_PROMPT

    for iteration in range(MAX_ITERATIONS):
        _truncate_tool_results_anthropic(messages, iteration + 1)
        logger.info("调用 API (anthropic), 第 %d 轮", iteration + 1)
        try:
            with client.messages.stream(
                model=MODEL,
                max_tokens=MAX_TOKENS,
                system=sys_prompt,
                messages=messages,
                tools=TOOLS,
            ) as stream:
                for event in stream:
                    if event.type == "content_block_delta":
                        delta = event.delta
                        if hasattr(delta, "text"):
                            yield {"event": "text_delta", "data": {"delta": delta.text}}
                final_message = stream.get_final_message()

            logger.info("API 返回, usage: %s", final_message.usage)
        except Exception as e:
            logger.error("API 调用失败: %s", e, exc_info=True)
            yield {"event": "error", "data": {"message": str(e)}}
            yield {"event": "done", "data": {}}
            return

        messages.append({"role": "assistant", "content": final_message.content})
        tool_blocks = [b for b in final_message.content if b.type == "tool_use"]

        if not tool_blocks:
            yield {"event": "done", "data": {}}
            return

        # 检查是否所有工具调用都是空参数
        all_empty = all(_check_tool_params(tb.name, tb.input) for tb in tool_blocks)
        if all_empty:
            empty_retries += 1
            if empty_retries >= _MAX_EMPTY_RETRIES:
                logger.warning("连续 %d 轮工具调用参数为空，终止循环", empty_retries)
                yield {"event": "done", "data": {}}
                return
        else:
            empty_retries = 0

        tool_results = []
        for tb in tool_blocks:
            param_err = _check_tool_params(tb.name, tb.input)
            if param_err:
                logger.warning("工具参数为空，跳过: %s", tb.name)
                tool_results.append({
                    "type": "tool_result", "tool_use_id": tb.id,
                    "content": param_err, "is_error": True,
                })
                continue

            yield {"event": "tool_start", "data": {"tool": tb.name, "input": tb.input}}
            tool_inp = dict(tb.input)
            if session_id and APP_MODE == "log-analyzer":
                tool_inp["_session_id"] = session_id
            result = exec_tool(tb.name, tool_inp)
            yield {"event": "tool_result", "data": {"tool": tb.name, "output": result[:MAX_DISPLAY_LEN]}}
            tool_results.append({
                "type": "tool_result", "tool_use_id": tb.id, "content": result,
            })

        messages.append({"role": "user", "content": tool_results})
        yield {"event": "thinking", "data": {}}

    logger.warning("达到最大迭代次数 %d", MAX_ITERATIONS)
    yield {"event": "error", "data": {"message": "达到最大迭代次数"}}
    yield {"event": "done", "data": {}}


def _run_openai_stream(messages: list, session_id: str = None):
    """OpenAI SDK 流式调用"""
    from openai import OpenAI
    client = OpenAI(base_url=BASE_URL, api_key=API_KEY)
    oai_tools = _tools_to_openai()
    empty_retries = 0

    # 构建 OpenAI 格式消息列表
    sys_prompt = build_system_prompt(session_id) if APP_MODE == "log-analyzer" else SYSTEM_PROMPT
    oai_messages = [{"role": "system", "content": sys_prompt}]
    for msg in messages:
        oai_messages.append({"role": msg["role"], "content": msg["content"]})

    for iteration in range(MAX_ITERATIONS):
        _truncate_tool_results_openai(oai_messages, iteration + 1)
        logger.info("调用 API (openai), 第 %d 轮", iteration + 1)
        try:
            stream = client.chat.completions.create(
                model=MODEL,
                max_tokens=MAX_TOKENS,
                messages=oai_messages,
                tools=oai_tools,
                stream=True,
            )

            content = ""
            tool_calls_acc = {}  # index -> {id, name, arguments}

            try:
                for chunk in stream:
                    choice = chunk.choices[0] if chunk.choices else None
                    if not choice or not choice.delta:
                        continue
                    delta = choice.delta

                    if delta.content:
                        yield {"event": "text_delta", "data": {"delta": delta.content}}
                        content += delta.content

                    if delta.tool_calls:
                        for tc in delta.tool_calls:
                            idx = tc.index
                            if idx not in tool_calls_acc:
                                tool_calls_acc[idx] = {"id": "", "name": "", "arguments": ""}
                            if tc.id:
                                tool_calls_acc[idx]["id"] = tc.id
                            if tc.function and tc.function.name:
                                tool_calls_acc[idx]["name"] = tc.function.name
                            if tc.function and tc.function.arguments:
                                tool_calls_acc[idx]["arguments"] += tc.function.arguments
            finally:
                stream.close()

            logger.info("API 返回, 第 %d 轮, 工具调用数: %d", iteration + 1, len(tool_calls_acc))

        except Exception as e:
            logger.error("API 调用失败: %s", e, exc_info=True)
            yield {"event": "error", "data": {"message": str(e)}}
            yield {"event": "done", "data": {}}
            return

        # 无工具调用，结束
        if not tool_calls_acc:
            yield {"event": "done", "data": {}}
            return

        # 解析参数并检查是否全部为空
        parsed_calls = []
        for idx, tc in sorted(tool_calls_acc.items()):
            # 兜底：部分代理不返回 tool call id
            if not tc["id"]:
                tc["id"] = f"call_{iteration}_{idx}"
            try:
                inp = json.loads(tc["arguments"]) if tc["arguments"] else {}
            except json.JSONDecodeError:
                inp = {}
            parsed_calls.append((tc, inp))

        all_empty = all(_check_tool_params(tc["name"], inp) for tc, inp in parsed_calls)
        if all_empty:
            empty_retries += 1
            if empty_retries >= _MAX_EMPTY_RETRIES:
                logger.warning("连续 %d 轮工具调用参数为空，终止循环", empty_retries)
                yield {"event": "done", "data": {}}
                return
        else:
            empty_retries = 0

        # 构建 assistant 消息（含 tool_calls）
        assistant_msg = {"role": "assistant", "content": content or None, "tool_calls": [
            {"id": tc["id"], "type": "function", "function": {"name": tc["name"], "arguments": tc["arguments"]}}
            for tc, _ in parsed_calls
        ]}
        oai_messages.append(assistant_msg)

        # 执行工具
        for tc, inp in parsed_calls:
            param_err = _check_tool_params(tc["name"], inp)
            if param_err:
                logger.warning("工具参数为空，跳过: %s", tc["name"])
                oai_messages.append({
                    "role": "tool", "tool_call_id": tc["id"], "content": param_err,
                })
                continue

            yield {"event": "tool_start", "data": {"tool": tc["name"], "input": inp}}
            tool_inp = dict(inp)
            if session_id and APP_MODE == "log-analyzer":
                tool_inp["_session_id"] = session_id
            result = exec_tool(tc["name"], tool_inp)
            yield {"event": "tool_result", "data": {"tool": tc["name"], "output": result[:MAX_DISPLAY_LEN]}}
            oai_messages.append({
                "role": "tool", "tool_call_id": tc["id"], "content": result,
            })

        yield {"event": "thinking", "data": {}}

    logger.warning("达到最大迭代次数 %d", MAX_ITERATIONS)
    yield {"event": "error", "data": {"message": "达到最大迭代次数"}}
    yield {"event": "done", "data": {}}


def run_agent_stream(messages: list, session_id: str = None):
    """Agent 循环：根据 API_FORMAT 选择对应 SDK 进行流式调用"""
    # 更新 session 活跃时间
    if session_id and _session_manager:
        _session_manager.touch_session(session_id)
    if API_FORMAT == "openai":
        yield from _run_openai_stream(messages, session_id=session_id)
    else:
        yield from _run_anthropic_stream(messages, session_id=session_id)
